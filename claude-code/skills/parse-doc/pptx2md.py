"""PPTX -> Markdown 변환 스크립트.

Usage: py pptx2md.py <input.pptx> <output.md>

이미지는 parsed/images/ 디렉터리에 추출되고, Markdown에 참조가 삽입됩니다.
"""
import os
import sys
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def shape_to_md(shape, slide_num, img_counter, images_dir, basename):
    lines = []
    images = []

    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        img_counter[0] += 1
        image = shape.image
        ext = image.content_type.split('/')[-1]
        if ext == 'jpeg':
            ext = 'jpg'
        img_name = f'{basename}_slide{slide_num}_img{img_counter[0]}.{ext}'
        img_path = os.path.join(images_dir, img_name)
        with open(img_path, 'wb') as f:
            f.write(image.blob)
        lines.append(f'![이미지](images/{img_name})')
        lines.append(f'<!-- Claude: 이 이미지를 Read로 읽어 내용을 해석하세요: {img_path} -->')
        lines.append('')
        images.append(img_path)
    elif shape.has_table:
        table = shape.table
        rows = []
        for row in table.rows:
            cells = [cell.text.replace('\n', ' ').strip() for cell in row.cells]
            rows.append(cells)
        if rows:
            lines.append('| ' + ' | '.join(rows[0]) + ' |')
            lines.append('| ' + ' | '.join(['---'] * len(rows[0])) + ' |')
            for row in rows[1:]:
                lines.append('| ' + ' | '.join(row) + ' |')
            lines.append('')
    elif shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            if para.level > 0:
                indent = '  ' * para.level
                lines.append(f'{indent}- {text}')
            else:
                lines.append(text)

    return lines, images


def pptx_to_md(input_path, output_path):
    prs = Presentation(input_path)
    md_lines = []
    all_images = []
    img_counter = [0]

    basename = os.path.splitext(os.path.basename(input_path))[0]
    output_dir = os.path.dirname(output_path)
    images_dir = os.path.join(output_dir, 'images')
    os.makedirs(images_dir, exist_ok=True)

    for i, slide in enumerate(prs.slides, 1):
        md_lines.append(f'## Slide {i}')
        md_lines.append('')

        title_text = None
        body_lines = []

        for shape in slide.shapes:
            is_title = (shape.has_text_frame
                        and slide.shapes.title is not None
                        and shape.shape_id == slide.shapes.title.shape_id)
            if is_title:
                title_text = shape.text_frame.text.strip()
            else:
                lines, images = shape_to_md(shape, i, img_counter, images_dir, basename)
                body_lines.extend(lines)
                all_images.extend(images)

        if title_text:
            md_lines.append(f'### {title_text}')
            md_lines.append('')

        md_lines.extend(body_lines)
        md_lines.append('')
        md_lines.append('---')
        md_lines.append('')

    with open(output_path, 'w', encoding='utf-8') as f:
        f.write('\n'.join(md_lines))

    print(f'Converted: {input_path} -> {output_path}')
    if all_images:
        print(f'Images extracted: {len(all_images)}개 -> {images_dir}/')
        for img in all_images:
            print(f'  - {os.path.basename(img)}')


if __name__ == '__main__':
    if len(sys.argv) != 3:
        print('Usage: py pptx2md.py <input.pptx> <output.md>')
        sys.exit(1)
    pptx_to_md(sys.argv[1], sys.argv[2])
